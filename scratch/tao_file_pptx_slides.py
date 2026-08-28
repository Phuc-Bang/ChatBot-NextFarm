#!/usr/bin/env python3
"""
tao_file_pptx_slides.py — Sinh file thuyet trinh PowerPoint (.pptx) tong quan du an NextFarm AI.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUT_PPTX = Path("d:/NextFarm/docs/SLIDES_BAO_CAO_NEXTFARM.pptx")
ROOT_PPTX = Path("d:/NextFarm/SLIDES_BAO_CAO_NEXTFARM.pptx")

# Modern Agricultural AI Color Palette
DARK_BG = RGBColor(15, 23, 42)       # Slate 900
LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
CARD_BG = RGBColor(255, 255, 255)
PRIMARY = RGBColor(22, 163, 74)      # Emerald 600
PRIMARY_LIGHT = RGBColor(236, 253, 245)
SECONDARY = RGBColor(37, 99, 235)    # Blue 600
AMBER = RGBColor(217, 119, 6)        # Amber 600
ROSE = RGBColor(225, 29, 72)         # Rose 600
TEXT_MAIN = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)
BORDER_COLOR = RGBColor(226, 232, 240)

def add_header(slide, title_text, category_text="NEXTFARM AI — BÀI TOÁN A: CHỐNG BỊA ĐẶT"):
    """Tao thanh header dong nhat cho cac slide."""
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = PRIMARY

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_MAIN

def add_card(slide, left, top, width, height, title, content_list, border_color=PRIMARY):
    """Tao hop noi dung (card) co border va text dep mat."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_MAIN
    p_title.space_after = Pt(6)

    for item in content_list:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(4)

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==================== SLIDE 1: COVER ====================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11), Inches(3.5))
    tf = tbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🌱 NEXTFARM AI ASSISTANT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "BÁO CÁO TỔNG KẾT KỸ THUẬT"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)

    p3 = tf.add_paragraph()
    p3.text = "Bài Toán A: Chống Bịa Đặt & Đảm Bảo An Toàn Trợ Lý Nông Nghiệp"
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.space_after = Pt(20)

    p4 = tf.add_paragraph()
    p4.text = "PoC Giai Đoạn 1 & 2 Hoàn Tất 100% · 367/367 Unit Tests PASSED · Tháng 08/2026"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = PRIMARY

    # ==================== SLIDE 2: 4 HIỆN TƯỢNG RỦI RO ====================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "1. Bối Cảnh & 4 Hiện Tượng Rủi Ro Cần Giải Quyết")

    add_card(s2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.5),
             "A1: Bịa Số Liệu Cảm Biến Vườn",
             ["LLM trần tự đoán độ ẩm 75%, 30°C dù không có cảm biến.",
              "Nguy cơ: Làm úng cây, cháy lá, chết hoa mùa vụ.",
              "✓ Giải pháp: Intent Router nhận diện và chặn ngay (0 token, 1ms)."],
             border_color=ROSE)

    add_card(s2, Inches(6.8), Inches(1.6), Inches(5.6), Inches(2.5),
             "A2: Bịa Tính Năng Ứng Dụng NextFarm",
             ["LLM tự nhận app có dự báo giá nông sản, kết nối cân...",
              "Nguy cơ: Gây hiểu lầm tính năng sản phẩm và khiếu nại dịch vụ.",
              "✓ Giải pháp: Phân loại product_feature và từ chối rõ ràng."],
             border_color=AMBER)

    add_card(s2, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.5),
             "A3: Sai Cây Trồng & Vùng Miền",
             ["Áp thời vụ miền Bắc cho miền Tây; tư vấn cây ngoài phạm vi.",
              "Nguy cơ: Sai lệch quy trình canh tác nghiêm trọng.",
              "✓ Giải pháp: Scope Check chặn cứng 3 cây; RRF cộng điểm vùng."],
             border_color=ROSE)

    add_card(s2, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.5),
             "A4: Tiếng Việt Không Dấu & Thổ Ngữ",
             ["Nông dân gõ 'ca chua bi sau', LLM đoán sai dấu hoặc nhầm 'kg'.",
              "Nguy cơ: Kê sai loại thuốc hoặc sai liều lượng phân bón.",
              "✓ Giải pháp: Chuẩn hoá 4 lớp + Trigram matching tầng CSDL."],
             border_color=PRIMARY)

    # ==================== SLIDE 3: KIẾN TRÚC 7 CHẶNG ====================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "2. Kiến Trúc Pipeline 7 Chặng & 3 Tầng Guardrail")

    steps_data = [
        ("1. Chuẩn Hoá (0ms)", "Unicode NFC, hạ chữ thường, từ điển đồng nghĩa nông học."),
        ("2. Intent Router (5ms)", "Tầng 1 Rule-based chặn 141 ca sớm; Tầng 2 LLM Few-shot Router."),
        ("3. Scope Check (4ms)", "Kiểm soát cứng 3 cây trồng: Lúa, Cà chua, Dưa chuột."),
        ("4. Hybrid RRF (12ms)", "FTS + Trigram + Vector Local; Cổng kiểm duyệt view indexable_chunk."),
        ("5. Evidence Pack", "Đóng gói JSON ngữ cảnh nguyên văn có chunk_id và Source Tier."),
        ("6. LLM Generator", "Gemini 3.1 Flash Lite sinh câu trả lời kèm mã trích dẫn [cid]."),
        ("7. Grounding Validator", "Rào chắn cuối: Kiểm tra JSON trích dẫn & đối soát số liệu tất định.")
    ]

    for idx, (st_t, st_d) in enumerate(steps_data):
        y_pos = 1.6 + (idx * 0.75)
        badge = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y_pos), Inches(0.5), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = PRIMARY
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.paragraphs[0].text = str(idx + 1)
        tf_b.paragraphs[0].font.size = Pt(13)
        tf_b.paragraphs[0].font.bold = True
        tf_b.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER

        t_card = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(y_pos), Inches(11.0), Inches(0.6))
        t_card.fill.solid()
        t_card.fill.fore_color.rgb = CARD_BG
        t_card.line.color.rgb = BORDER_COLOR
        tf_c = t_card.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        
        r1 = p_c.add_run()
        r1.text = st_t + ": "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN

        r2 = p_c.add_run()
        r2.text = st_d
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # ==================== SLIDE 4: KHO TRI THỨC ====================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "3. Quản Trị Kho Tri Thức & Quy Trình Kiểm Duyệt 2 Luồng")

    add_card(s4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.5),
             "Tài Liệu Khuyến Nông (18/31 Đã Duyệt)",
             ["Thu thập chính thống từ các cổng thông tin .gov.vn.",
              "18 tài liệu hướng dẫn kỹ thuật chuẩn đã được phê duyệt.",
              "13 tài liệu dạng tin bài tổng hợp/tin hoạt động bị loại bỏ."],
             border_color=PRIMARY)

    add_card(s4, Inches(6.8), Inches(1.6), Inches(5.6), Inches(2.5),
             "Kho RAG Indexable (185/292 Chunks)",
             ["185 chunks kỹ thuật chuẩn mở khóa vào kho tìm kiếm.",
              "Lúa: 87 chunks · Dưa chuột: 62 chunks · Cà chua: 36 chunks.",
              "Khóa chặt bằng mã băm SHA-256 (DEC-005)."],
             border_color=PRIMARY)

    add_card(s4, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.5),
             "Chunk Rủi Ro Cao (31 Chunk Duyệt Lẻ)",
             ["Chứa liều lượng phân bón, thuốc BVTV, nồng độ hoạt chất.",
              "24 chunk đạt chuẩn được cấp phép; 7 chunk bị loại bỏ.",
              "Trạng thái: 100% chunk rủi ro đã được xử lý (0 chunk chờ)."],
             border_color=AMBER)

    add_card(s4, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.5),
             "Dữ Liệu Ground Truth (65/141 Facts)",
             ["65 fact số liệu chuẩn về độ pH, lượng giống, đợt bón thúc.",
              "Xác thực thủ công làm ground truth để kiểm thử tự động.",
              "Grounding Validator dùng fact này để so khớp tất định."],
             border_color=SECONDARY)

    # ==================== SLIDE 5: 3 MỤC §40.2 ====================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "4. Tối Ưu Hoá & Giải Quyết 3 Tham Số Kỹ Thuật §40.2")

    add_card(s5, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.2),
             "Mục 4: Sweep Kích Thước Chunk",
             ["Thực nghiệm 4 dải: 800, 1000, 1200, 1500 ký tự.",
              "Đo trên 18 tài liệu & 141 facts.",
              "Kết quả: 1.200 ký tự đạt 100% Fact Coverage & 100% Hit Rate.",
              "Quyết định: Giữ chuẩn 1.200 ký tự, không làm rách bảng phân bón.",
              "Báo cáo: P16_chunk_size_sweep.md"],
             border_color=PRIMARY)

    add_card(s5, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.2),
             "Mục 9: LLM Few-Shot Router",
             ["Bổ sung tầng LLMFewShotRouter khi rule-based không chắc chắn.",
              "Tích hợp 35+ ví dụ mẫu chuẩn hóa theo 6 intent.",
              "Áp dụng quy tắc thiên lệch an toàn §11.4.",
              "Tính xác suất tin cậy softmax thực (0.85 - 0.95).",
              "Bảo đảm 100% an toàn cho câu hỏi lệnh thiết bị."],
             border_color=SECONDARY)

    add_card(s5, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.2),
             "Mục 11: Trọng Số Vùng Miền",
             ["Xây dựng bộ 30 test case đặc thù vùng miền.",
              "Quét dải hệ số ưu tiên từ 0.00 đến 0.30.",
              "Kết quả: he_so = 0.10 nâng Top-1 khớp vùng lên 63.3%, Top-3 đạt 73.3%.",
              "Không làm méo mó xếp hạng tài liệu liên quan.",
              "Báo cáo: P17_region_weighting_sweep.md"],
             border_color=AMBER)

    # ==================== SLIDE 6: ĐỐI CHỨNG C0 - C1 - C2 ====================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "5. Bảng Đối Chứng Định Lượng (C0 · C1 · C2)")

    rows = 7
    cols = 4
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(11.7)
    height = Inches(5.0)

    table_shape = s6.shapes.add_table(rows, cols, left, top, width, height)
    t = table_shape.table
    t.columns[0].width = Inches(3.8)
    t.columns[1].width = Inches(2.6)
    t.columns[2].width = Inches(2.6)
    t.columns[3].width = Inches(2.7)

    headers = ["Chỉ Số Đo Lường (222 Test Cases)", "C0: LLM Trần", "C1: RAG Thuần", "C2: NextFarm PoC"]
    for i, h in enumerate(headers):
        cell = t.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY if i == 3 else DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

    data = [
        ("Số ca bịa số liệu nông học", "61 ca (27.5%)", "23 ca (10.4%)", "0 ca (0.0% Triệt tiêu)"),
        ("Tỷ lệ trả lời sai (False Answer Rate)", "77.0%", "23.0%", "0.9%"),
        ("Độ chính xác khi trả lời", "1.2%", "23.9%", "90.9%"),
        ("Bắt ca cần từ chối (Abstention Recall)", "3.4%", "69.6%", "100.0% (148/148 ca)"),
        ("Độ trễ trung vị (p50)", "2.621 ms", "2.555 ms", "15 ms (Giảm 99.4%)"),
        ("Chi phí API / 222 câu", "$0.0369", "$0.1231", "$0.0527 (~$0.0002/câu)")
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, val in enumerate(row_data):
            cell = t.cell(row_idx, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_LIGHT if col_idx == 3 else (RGBColor(241, 245, 249) if row_idx % 2 == 0 else CARD_BG)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_MAIN
            if col_idx == 3 or col_idx == 0:
                p.font.bold = True

    # ==================== SLIDE 7: CHẤM CHUYÊN GIA ====================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "6. Kết Quả Chấm Chuyên Gia C2 (50 Câu Thực Tế)")

    add_card(s7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
             "Điểm Sáng & Chất Lượng Thực Tế",
             ["Chấm thật bởi Đội kỹ thuật NextFarm.",
              "Đọc nguyên văn câu trả lời thực tế của hệ thống.",
              "Trích dẫn đúng 100% tài liệu gốc cho câu trả lời chuẩn (pH 6.0-6.5, ngâm hạt 40-50°C).",
              "Từ chối bẫy kê đơn thuốc BVTV và bẫy bò lan chính xác.",
              "Điểm số phân hóa tự nhiên theo 8+ tổ hợp thực tế."],
             border_color=PRIMARY)

    add_card(s7, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.2),
             "Nhận Diện Khuyết Điểm Thực Tế",
             ["Phát hiện lỗi tại Câu 41 (hỏi độ ẩm) và Câu 42 (hỏi dữ liệu vườn):",
              "Hệ thống bị Scope Check chặn nhầm sang template hỏi lại cây trồng ('can_lam_ro').",
              "Khuyết điểm: Từ chối an toàn nhưng chưa giải thích đúng lý do thiếu cảm biến vườn.",
              "Hạ điểm chuẩn xác (C1:3, C2:2, C3:2) để phản ánh trung thực thực trạng."],
             border_color=AMBER)

    # ==================== SLIDE 8: PRODUCTION & DEVOPS ====================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "7. Đóng Gói Sản Xuất & Đảm Bảo Chất Lượng (Production Ready)")

    add_card(s8, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.5),
             "Đóng Gói Docker Compose 1 Lệnh",
             ["Dockerfile Python 3.11-slim tối ưu, non-root user nextfarm.",
              "docker-compose.prod.yml gồm 2 container: App + DB pgvector.",
              "Tích hợp healthcheck và network cô lập an toàn."],
             border_color=PRIMARY)

    add_card(s8, Inches(6.8), Inches(1.6), Inches(5.6), Inches(2.5),
             "Bảo Mật & Local Embedding",
             ["Embedding chạy 100% On-Premise, không gửi kho tri thức ra ngoài.",
              "Cửa kiểm soát quyền hạn /admin mặc định an toàn.",
              "Biến môi trường và API keys được cô lập tuyệt đối."],
             border_color=PRIMARY)

    add_card(s8, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.5),
             "Bộ Kiểm Thử Tự Động (367 Tests Xanh)",
             ["367 Unit Tests kiểm định toàn bộ pipeline.",
              "Phủ từ Chuẩn hóa, Intent Router, Scope, Retrieval đến Grounding.",
              "Tự động chạy và đồng bộ trên CI/CD."],
             border_color=PRIMARY)

    add_card(s8, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.5),
             "4 Giao Diện Web Đa Dạng",
             ["Trang Chat nông dân (/) · Trang Quản trị (/admin).",
              "Trang Báo cáo kỹ thuật (/report) · Trang Chấm điểm (/expert).",
              "Hỗ trợ giao diện Sáng / Tối và thiết kế đáp ứng (Responsive)."],
             border_color=SECONDARY)

    # ==================== SLIDE 9: GIỚI HẠN & LỘ TRÌNH ====================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "8. Giới Hạn Đã Biết & Lộ Trình Mở Rộng Tiếp Theo")

    add_card(s9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
             "Giới Hạn Đã Biết (Trung Thực & Minh Bạch)",
             ["1. Chưa kết nối API cảm biến IoT vườn thật: PoC Bài toán A tập trung chống bịa đặt nên tạm thời từ chối câu hỏi dữ liệu vườn.",
              "2. Chưa có tài liệu sản phẩm riêng NextFarm: Câu hỏi bảng giá thiết bị đang được phân loại vào product_feature.",
              "3. Fine-tuning self-host (P12): Cần GPU chuyên dụng (≥16GB VRAM) nếu muốn tự host hoàn toàn LLM thay vì dùng API.",
              "4. Kho tri thức hiện tại tập trung vào 3 cây trồng chủ lực: Lúa, Cà chua, Dưa chuột (185 chunks indexable)."],
             border_color=AMBER)

    add_card(s9, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.2),
             "Lộ Trình Phát Triển Giai Đoạn Tiếp Theo",
             ["Giai Đoạn 1.5 (Hiện Tại):",
              "• Đóng gói Docker, bàn giao mã nguồn & báo cáo tổng kết.",
              "Giai Đoạn 2 (Mở Rộng Canh Tác):",
              "• Mở rộng kho tri thức cho cây ăn quả/công nghiệp (Sầu riêng, Xoài, Cà phê, Hồ tiêu...) lên 2.000+ chunks.",
              "• Tích hợp cẩm nang sản phẩm NextFarm.",
              "Giai Đoạn 3 (Hệ Sinh Thái IoT):",
              "• Kết nối API trạm đo thực tế theo từng vườn (Bài toán B).",
              "• Thiết lập xác thực 2 lớp (2FA) cho lệnh điều khiển bơm/van."],
             border_color=PRIMARY)

    # ==================== SLIDE 10: KẾT LUẬN ====================
    s10 = prs.slides.add_slide(blank_layout)
    bg10 = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg10.fill.solid()
    bg10.fill.fore_color.rgb = DARK_BG
    bg10.line.fill.background()

    tbox10 = s10.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(4.0))
    tf10 = tbox10.text_frame
    tf10.word_wrap = True

    p = tf10.paragraphs[0]
    p.text = "TỔNG KẾT & BÀN GIAO DỰ ÁN"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(14)

    bullets = [
        "✓ Triệt tiêu 100% hiện tượng bịa số liệu nông học (61 ca -> 0 ca).",
        "✓ Bắt đúng 100% ca cần từ chối (148/148 ca bẫy và ngoài phạm vi).",
        "✓ Tốc độ xử lý p50 đạt 15ms (giảm 99.4% thời gian phản hồi).",
        "✓ Đóng gói Docker Compose sản xuất hoàn chỉnh 1 lệnh.",
        "✓ Toàn bộ 367 Unit Tests tự động xanh 100% và sẵn sàng vận hành."
    ]

    for b in bullets:
        p_b = tf10.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(16)
        p_b.font.color.rgb = RGBColor(255, 255, 255)
        p_b.space_after = Pt(8)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    prs.save(str(ROOT_PPTX))
    print(f"[OK] Đã tạo thành công file PowerPoint tại:\n- {OUT_PPTX}\n- {ROOT_PPTX}")

if __name__ == "__main__":
    main()
